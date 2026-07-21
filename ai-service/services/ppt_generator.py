"""
ppt_generator.py - Membuat PPTX opsional untuk guru dari materi PDF.

Alur:
- Provider LLM membuat outline slide terstruktur.
- python-pptx mengisi outline ke template PPTX yang sudah disediakan.
- File PPTX dikembalikan sebagai bytes ke backend-api untuk disimpan dan
  disediakan sebagai download guru.
"""

import io
import json
import logging
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Pt

from services.nim_client import NIMAPIError, _parse_llm_json, _get_config
import requests
import time

logger = logging.getLogger("ai-service.ppt_generator")

DEFAULT_TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(__file__)),
    "templates",
    "template-pendidikan.pptx",
)


def _remove_slide(prs: Presentation, slide):
    slide_id_list = prs.slides._sldIdLst  # python-pptx tidak punya public delete API.
    slide_id = slide.slide_id
    for index, sld_id in enumerate(slide_id_list):
        if sld_id.id == slide_id:
            r_id = sld_id.rId
            prs.part.drop_rel(r_id)
            del slide_id_list[index]
            return


def _text_shapes(slide):
    return [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text_frame
        and shape.text_frame.text.strip()
    ]


def _set_shape_text(shape, text: str, font_size: Optional[int] = None):
    shape.text = text
    if font_size:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)


def _fill_template_slide(slide, title: str, bullets: list[str]):
    shapes = _text_shapes(slide)
    if not shapes:
        return
    _set_shape_text(shapes[0], title.upper(), 28)
    if len(shapes) >= 2:
        bullet_text = "\n".join(f"• {bullet}" for bullet in bullets)
        _set_shape_text(shapes[1], bullet_text, 18)


def _build_outline_prompt(title: str, materi_text: str) -> str:
    return (
        "Buat outline presentasi PowerPoint untuk guru berdasarkan materi berikut. "
        "Tujuannya untuk menjelaskan rangkuman, poin-poin penting, contoh, dan penutup. "
        "Jawab HANYA JSON object valid dengan format: "
        '{"slides":[{"title":"...","bullets":["...","..."]}]}. '
        "Buat 6-8 slide. Slide pertama harus judul dan tujuan belajar. "
        "Slide terakhir harus ringkasan singkat. Setiap slide maksimal 5 bullet, "
        "setiap bullet maksimal 14 kata. Bahasa Indonesia, jelas, dan cocok untuk kelas.\n\n"
        f"Judul materi: {title}\n\nMateri:\n{materi_text}"
    )


def generate_ppt_outline(title: str, materi_text: str) -> list[dict]:
    config = _get_config()
    prompt = _build_outline_prompt(title, materi_text)
    if config["provider"] == "nvidia":
        return _generate_ppt_outline_nvidia(config, prompt)

    return _generate_ppt_outline_gemini(config, prompt)


def _parse_outline_content(content_str: str) -> list[dict]:
    parsed = _parse_llm_json(content_str, "ppt").get("parsed")
    if isinstance(parsed, list):
        slides = parsed
    elif isinstance(parsed, dict) and isinstance(parsed.get("slides"), list):
        slides = parsed["slides"]
    else:
        raise NIMAPIError("Response LLM untuk PPT tidak berisi slides JSON valid.")
    return _normalize_slides(slides)


def _generate_ppt_outline_nvidia(config: dict, prompt: str) -> list[dict]:
    payload = {
        "model": config["model"],
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.2,
        "top_p": 0.7,
        "max_tokens": 1800,
        "stream": False,
    }
    headers = {
        "Authorization": f"Bearer {config['api_key']}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    }
    url = f"{config['base_url'].rstrip('/')}/chat/completions"

    last_error = None
    for attempt in range(1, config["max_retries"] + 2):
        try:
            response = requests.post(url, headers=headers, json=payload, timeout=config["timeout"])
            if response.status_code == 429:
                last_error = NIMAPIError("Rate limit tercapai pada NVIDIA NIM API (429).")
                time.sleep(max(config["rate_limit_sleep"], min(2 ** attempt, 8)))
                continue

            response.raise_for_status()
            data = response.json()
            return _parse_outline_content(data["choices"][0]["message"]["content"])
        except requests.exceptions.HTTPError as e:
            status_code = e.response.status_code if e.response is not None else 500
            safe_error = str(e).replace(config["api_key"], "[redacted]")
            last_error = NIMAPIError(f"Gagal menghubungi NVIDIA NIM API (status {status_code}): {safe_error}")
            if status_code == 429 or status_code >= 500:
                wait_seconds = max(config["rate_limit_sleep"], min(2 ** attempt, 8)) if status_code == 429 else min(2 ** attempt, 8)
                time.sleep(wait_seconds)
                continue
            raise last_error from None
        except requests.exceptions.RequestException as e:
            safe_error = str(e).replace(config["api_key"], "[redacted]")
            last_error = NIMAPIError(f"Gagal menghubungi NVIDIA NIM API: {safe_error}")
            time.sleep(min(2 ** attempt, 8))
        except (KeyError, IndexError, json.JSONDecodeError) as e:
            last_error = NIMAPIError(f"Response NVIDIA NIM untuk PPT tidak valid: {e}")
            continue
        except NIMAPIError as e:
            logger.warning("NVIDIA NIM API gagal membuat PPT (percobaan %s): %s", attempt, e)
            last_error = e
            continue

    raise (last_error or NIMAPIError("Gagal membuat outline PPT setelah beberapa percobaan.")) from None


def _generate_ppt_outline_gemini(config: dict, prompt: str) -> list[dict]:
    payload = {
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.25,
            "maxOutputTokens": 1800,
            "responseMimeType": "application/json",
            "thinkingConfig": {"thinkingBudget": 0},
        },
    }

    url = f"{config['base_url']}/models/{config['model']}:generateContent"
    params = {"key": config["api_key"]}
    headers = {"Content-Type": "application/json", "Accept": "application/json"}

    last_error = None
    for attempt in range(1, config["max_retries"] + 2):
        try:
            response = requests.post(url, headers=headers, params=params, json=payload, timeout=config["timeout"])
            if response.status_code == 429:
                last_error = NIMAPIError("Rate limit tercapai pada Gemini API (429).")
                time.sleep(max(config["rate_limit_sleep"], min(2 ** attempt, 8)))
                continue

            response.raise_for_status()
            data = response.json()
            content_str = "".join(
                part.get("text", "")
                for part in data["candidates"][0]["content"].get("parts", [])
            )
            return _parse_outline_content(content_str)
        except requests.exceptions.HTTPError as e:
            status_code = e.response.status_code if e.response is not None else 500
            safe_error = str(e).replace(config["api_key"], "[redacted]")
            last_error = NIMAPIError(f"Gagal menghubungi Gemini API (status {status_code}): {safe_error}")
            if status_code == 429 or status_code >= 500:
                wait_seconds = max(config["rate_limit_sleep"], min(2 ** attempt, 8)) if status_code == 429 else min(2 ** attempt, 8)
                time.sleep(wait_seconds)
                continue
            raise last_error from None
        except requests.exceptions.RequestException as e:
            safe_error = str(e).replace(config["api_key"], "[redacted]")
            last_error = NIMAPIError(f"Gagal menghubungi Gemini API: {safe_error}")
            time.sleep(min(2 ** attempt, 8))
        except (KeyError, IndexError, json.JSONDecodeError) as e:
            last_error = NIMAPIError(f"Response Gemini untuk PPT tidak valid: {e}")
            continue
        except NIMAPIError as e:
            logger.warning("Gemini API gagal membuat PPT (percobaan %s): %s", attempt, e)
            last_error = e
            continue

    raise (last_error or NIMAPIError("Gagal membuat outline PPT setelah beberapa percobaan.")) from None


def _normalize_slides(slides: list[dict]) -> list[dict]:
    normalized = []
    for slide in slides[:8]:
        title = str(slide.get("title") or "").strip()
        bullets = slide.get("bullets") if isinstance(slide.get("bullets"), list) else []
        clean_bullets = [str(item).strip() for item in bullets if str(item).strip()][:5]
        if title and clean_bullets:
            normalized.append({"title": title[:90], "bullets": clean_bullets})
    if len(normalized) < 2:
        raise NIMAPIError("Outline PPT terlalu sedikit untuk dibuat menjadi presentasi.")
    return normalized


def build_pptx(title: str, outline: list[dict], template_path: Optional[str] = None) -> bytes:
    template = template_path or os.getenv("PPT_TEMPLATE_PATH", DEFAULT_TEMPLATE_PATH)
    prs = Presentation(template)

    # Template berisi text box biasa, bukan placeholder. Isi langsung slide
    # template supaya desain/ilustrasi asli tetap utuh.
    cover = prs.slides[0]
    cover_shapes = _text_shapes(cover)
    if cover_shapes:
        _set_shape_text(cover_shapes[0], title.upper(), 34)
    if len(cover_shapes) >= 2:
        _set_shape_text(cover_shapes[1], "RANGKUMAN DAN POIN-POIN PENTING", 16)

    max_content_slides = max(0, len(prs.slides) - 1)
    selected_outline = outline[:max_content_slides]
    for index, slide_data in enumerate(selected_outline, start=1):
        _fill_template_slide(prs.slides[index], slide_data["title"], slide_data["bullets"])

    keep_count = 1 + len(selected_outline)
    while len(prs.slides) > keep_count:
        _remove_slide(prs, prs.slides[-1])

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def generate_pptx(title: str, materi_text: str) -> bytes:
    outline = generate_ppt_outline(title, materi_text)
    return build_pptx(title, outline)
